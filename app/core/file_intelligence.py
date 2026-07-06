import os
import io
import csv
import json
import zipfile
from typing import Optional

# Graceful optional imports for document libraries
try:
    import pypdf
    HAS_PYPDF = True
except ImportError:
    HAS_PYPDF = False

try:
    import docx2txt
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    import openpyxl
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False

try:
    import pptx
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


def parse_binary_document(filename: str, file_bytes: bytes) -> str:
    """Parse Excel, PDF, Word, PPTX, or ZIP documents from raw bytes."""
    ext = os.path.splitext(filename.lower())[1]

    if ext == ".pdf":
        if not HAS_PYPDF:
            return "[Error: pypdf is not installed. Run 'pip install pypdf' to enable PDF ingestion.]"
        try:
            reader = pypdf.PdfReader(io.BytesIO(file_bytes))
            text_parts = []
            for i, page in enumerate(reader.pages):
                t = page.extract_text()
                if t:
                    text_parts.append(f"--- Page {i+1} ---\n{t}")
            return "\n\n".join(text_parts)
        except Exception as e:
            return f"[Error parsing PDF '{filename}': {e}]"

    elif ext in [".docx", ".doc"]:
        if not HAS_DOCX:
            return "[Error: docx2txt is not installed. Run 'pip install docx2txt' to enable Word document ingestion.]"
        try:
            return docx2txt.process(io.BytesIO(file_bytes))
        except Exception as e:
            return f"[Error parsing Word doc '{filename}': {e}]"

    elif ext in [".xlsx", ".xls"]:
        if not HAS_OPENPYXL:
            return "[Error: openpyxl is not installed. Run 'pip install openpyxl' to enable Excel sheet ingestion.]"
        try:
            wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
            sheets_text = []
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                rows_text = []
                for row in sheet.iter_rows(values_only=True):
                    if any(row):  # skip completely empty rows
                        row_str = " | ".join(str(cell) if cell is not None else "" for cell in row)
                        rows_text.append(row_str)
                sheets_text.append(f"--- Sheet: {sheet_name} ---\n" + "\n".join(rows_text))
            return "\n\n".join(sheets_text)
        except Exception as e:
            return f"[Error parsing Excel sheet '{filename}': {e}]"

    elif ext == ".pptx":
        if not HAS_PPTX:
            return "[Error: python-pptx is not installed. Run 'pip install python-pptx' to enable PowerPoint ingestion.]"
        try:
            prs = pptx.Presentation(io.BytesIO(file_bytes))
            slides_text = []
            for i, slide in enumerate(prs.slides):
                slide_parts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_parts.append(shape.text.strip())
                slides_text.append(f"--- Slide {i+1} ---\n" + "\n".join(slide_parts))
            return "\n\n".join(slides_text)
        except Exception as e:
            return f"[Error parsing PowerPoint '{filename}': {e}]"

    elif ext == ".zip":
        try:
            with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
                names = z.namelist()
                structure = "\n".join(f"- {name}" for name in names)
                file_contents = []
                # Read first 5 text files as a sample
                text_count = 0
                for name in names:
                    if name.endswith("/") or any(x in name.lower() for x in ["__pycache__", ".git", ".idea", "node_modules"]):
                        continue
                    # Read text-based formats
                    sub_ext = os.path.splitext(name.lower())[1]
                    if sub_ext in [".txt", ".py", ".js", ".json", ".md", ".html", ".css", ".xml", ".yaml", ".csv", ".log"]:
                        try:
                            with z.open(name) as f:
                                content = f.read().decode("utf-8", errors="ignore")
                                file_contents.append(f"--- Zip File: {name} ---\n{content[:4000]}")
                                text_count += 1
                                if text_count >= 5:
                                    break
                        except Exception:
                            pass
                
                parts = [f"ZIP Archive contains {len(names)} files.\nStructure:\n{structure}"]
                if file_contents:
                    parts.append("\n\nSample Contents:\n" + "\n\n".join(file_contents))
                return "\n\n".join(parts)
        except Exception as e:
            return f"[Error parsing ZIP archive '{filename}': {e}]"

    # Default text decode fallback
    try:
        return file_bytes.decode("utf-8", errors="ignore")
    except Exception as e:
        return f"[Unsupported binary file type '{filename}': {e}]"
